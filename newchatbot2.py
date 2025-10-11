import tkinter as tk
from tkinter import ttk, scrolledtext, filedialog, messagebox
import threading
import time
import os
from datetime import datetime
import re
import math
import random
from collections import Counter, defaultdict

# For file processing
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_WRITE_AVAILABLE = True
except ImportError:
    PPTX_WRITE_AVAILABLE = False

class AdvancedNLG:
    """Advanced Natural Language Generator - Creates human-like responses"""
    
    def __init__(self):
        # Enhanced synonym database
        self.synonyms = {
            'show': ['demonstrate', 'illustrate', 'reveal', 'indicate', 'display', 'exhibit'],
            'explain': ['describe', 'clarify', 'elucidate', 'elaborate', 'detail', 'outline'],
            'important': ['significant', 'crucial', 'vital', 'essential', 'critical', 'key'],
            'use': ['utilize', 'employ', 'apply', 'implement', 'adopt'],
            'help': ['assist', 'aid', 'support', 'facilitate', 'enable'],
            'make': ['create', 'produce', 'generate', 'construct', 'build'],
            'get': ['obtain', 'acquire', 'receive', 'gain', 'secure'],
            'give': ['provide', 'offer', 'present', 'supply', 'deliver'],
            'think': ['believe', 'consider', 'feel', 'suppose', 'assume'],
            'know': ['understand', 'comprehend', 'recognize', 'realize', 'grasp'],
            'see': ['observe', 'notice', 'perceive', 'detect', 'witness'],
            'find': ['discover', 'locate', 'identify', 'detect', 'uncover'],
            'allow': ['enable', 'permit', 'facilitate', 'authorize', 'let'],
            'include': ['contain', 'comprise', 'encompass', 'incorporate', 'feature'],
        }
        
        # Conversational intros (more natural)
        self.intros = [
            "Based on what I found in the document,",
            "From what I can gather,",
            "According to the content,",
            "Here's what I understand:",
            "Let me explain what I found:",
            "The document suggests that",
            "From my analysis,",
            "What I'm seeing here is that",
            "It appears that",
            "The information indicates that",
        ]
        
        # Natural transitions
        self.transitions = {
            'addition': ['Also,', 'Additionally,', 'Furthermore,', 'Moreover,', 'Plus,', 'What\'s more,'],
            'elaboration': ['In fact,', 'Specifically,', 'More precisely,', 'To be exact,', 'In particular,'],
            'contrast': ['However,', 'On the other hand,', 'That said,', 'Nevertheless,', 'But,'],
            'example': ['For example,', 'For instance,', 'Take this case:', 'Consider:', 'Like:'],
            'result': ['Therefore,', 'As a result,', 'Consequently,', 'Thus,', 'So,'],
            'emphasis': ['Indeed,', 'In fact,', 'Actually,', 'Importantly,', 'Notably,'],
        }
        
        # Sentence starters for variety
        self.sentence_starters = [
            "It's worth noting that", "Interestingly,", "What's important here is that",
            "I should mention that", "One key point is that", "It's clear that",
            "Essentially,", "In essence,", "Basically,", "Simply put,"
        ]
        
        # Ending phrases for engagement
        self.engagers = [
            "Does this help clarify things?",
            "Would you like me to elaborate on any part?",
            "Is there a specific aspect you'd like me to dive deeper into?",
            "Let me know if you need more details!",
            "Feel free to ask if something's unclear.",
            "I'm happy to explain further if needed.",
        ]
    
    def create_human_paragraph(self, sentences, style='informative'):
        """Create a natural, flowing paragraph"""
        if not sentences:
            return ""
        
        # Start with an intro
        result = [random.choice(self.intros) + " " + sentences[0]]
        
        for i in range(1, min(len(sentences), 4)):
            # Add variety with transitions
            if i == 1 and random.random() > 0.3:
                trans_type = random.choice(['elaboration', 'addition'])
                result.append(random.choice(self.transitions[trans_type]) + " " + sentences[i])
            elif i == 2 and random.random() > 0.4:
                result.append(random.choice(self.sentence_starters) + " " + sentences[i].lower())
            else:
                result.append(sentences[i])
        
        # Join naturally
        paragraph = " ".join(result)
        
        # Add engagement if appropriate
        if style == 'helpful' and random.random() > 0.5:
            paragraph += "\n\n" + random.choice(self.engagers)
        
        return paragraph
    
    def paraphrase_intelligently(self, text):
        """Intelligent paraphrasing using multiple techniques"""
        # Split into words
        words = text.split()
        result = []
        
        for i, word in enumerate(words):
            word_clean = word.lower().strip('.,!?;:')
            
            # Replace with synonym occasionally
            if word_clean in self.synonyms and random.random() > 0.6:
                synonym = random.choice(self.synonyms[word_clean])
                # Preserve capitalization
                if word[0].isupper():
                    synonym = synonym.capitalize()
                # Preserve punctuation
                if word[-1] in '.,!?;:':
                    synonym += word[-1]
                result.append(synonym)
            else:
                result.append(word)
        
        return ' '.join(result)
    
    def restructure_for_naturalness(self, sentence):
        """Make sentences sound more conversational"""
        # Remove overly formal patterns
        sentence = re.sub(r'^The document states that', 'Basically,', sentence)
        sentence = re.sub(r'^It is important to note that', 'Keep in mind that', sentence)
        sentence = re.sub(r'^It should be noted that', 'Worth mentioning -', sentence)
        
        return sentence

class ContextualUnderstanding:
    """Understands question context and intent deeply"""
    
    def __init__(self):
        self.question_patterns = {
            'definition': {
                'patterns': [r'what is', r'what are', r'define', r'meaning of', r'definition of'],
                'intent': 'seeking_definition'
            },
            'explanation': {
                'patterns': [r'why', r'how does', r'explain', r'reason', r'cause'],
                'intent': 'seeking_explanation'
            },
            'procedure': {
                'patterns': [r'how to', r'steps', r'process', r'method', r'way to'],
                'intent': 'seeking_steps'
            },
            'comparison': {
                'patterns': [r'difference', r'compare', r'versus', r'vs', r'better than'],
                'intent': 'seeking_comparison'
            },
            'examples': {
                'patterns': [r'example', r'instance', r'such as', r'like what'],
                'intent': 'seeking_examples'
            },
            'listing': {
                'patterns': [r'list', r'types of', r'kinds of', r'categories'],
                'intent': 'seeking_list'
            },
            'yes_no': {
                'patterns': [r'^is ', r'^are ', r'^does ', r'^do ', r'^can ', r'^will '],
                'intent': 'seeking_confirmation'
            },
        }
        
        self.stop_words = {
            'the', 'is', 'at', 'which', 'on', 'a', 'an', 'as', 'are', 'was', 'were',
            'been', 'be', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would',
            'should', 'could', 'may', 'might', 'must', 'can', 'of', 'for', 'to', 'in',
            'by', 'with', 'from', 'about', 'into', 'through', 'during', 'before', 'after'
        }
    
    def analyze_question(self, question):
        """Deep analysis of question intent"""
        question_lower = question.lower()
        
        # Detect question type
        q_type = 'general'
        intent = 'seeking_information'
        
        for qtype, data in self.question_patterns.items():
            for pattern in data['patterns']:
                if re.search(pattern, question_lower):
                    q_type = qtype
                    intent = data['intent']
                    break
        
        # Extract key concepts (not just keywords)
        concepts = self.extract_concepts(question)
        
        # Determine complexity
        complexity = 'simple' if len(question.split()) < 8 else 'complex'
        
        return {
            'type': q_type,
            'intent': intent,
            'concepts': concepts,
            'complexity': complexity,
            'requires_detail': 'explain' in question_lower or 'detail' in question_lower
        }
    
    def extract_concepts(self, text):
        """Extract main concepts from question"""
        words = re.findall(r'\b[a-z]{3,}\b', text.lower())
        concepts = [w for w in words if w not in self.stop_words]
        
        # Find phrases (2-3 word combinations)
        words_list = text.lower().split()
        phrases = []
        for i in range(len(words_list) - 1):
            if words_list[i] not in self.stop_words and words_list[i+1] not in self.stop_words:
                phrases.append(f"{words_list[i]} {words_list[i+1]}")
        
        return {
            'keywords': concepts[:5],
            'phrases': phrases[:3]
        }

class SemanticMatcher:
    """Matches questions to content semantically"""
    
    def __init__(self):
        self.context_understanding = ContextualUnderstanding()
    
    def compute_semantic_score(self, query_concepts, sentence):
        """Compute how relevant a sentence is"""
        sentence_lower = sentence.lower()
        score = 0
        matches = []
        
        # Check keywords
        for keyword in query_concepts['keywords']:
            if keyword in sentence_lower:
                score += 2
                matches.append(keyword)
                # Bonus for multiple occurrences
                score += sentence_lower.count(keyword) * 0.5
        
        # Check phrases (worth more)
        for phrase in query_concepts['phrases']:
            if phrase in sentence_lower:
                score += 5
                matches.append(phrase)
        
        # Length consideration (not too short, not too long)
        word_count = len(sentence.split())
        if 10 < word_count < 50:
            score += 1
        
        return {
            'score': score,
            'matches': matches,
            'text': sentence
        }
    
    def find_relevant_content(self, question, content, top_n=5):
        """Find most relevant content for the question"""
        analysis = self.context_understanding.analyze_question(question)
        concepts = analysis['concepts']
        
        # Split content into sentences
        sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 15]
        
        if not sentences:
            return []
        
        # Score each sentence
        scored = []
        for i, sent in enumerate(sentences):
            result = self.compute_semantic_score(concepts, sent)
            # Add position bonus for first 20 sentences
            if i < 20:
                result['score'] += (20 - i) * 0.1
            scored.append(result)
        
        # Sort by score
        scored.sort(key=lambda x: x['score'], reverse=True)
        
        # Filter out low scores and return top N
        relevant = [s for s in scored if s['score'] > 0]
        return relevant[:top_n], analysis

class AdvancedResponseEngine:
    """Generates highly natural, context-aware responses"""
    
    def __init__(self):
        self.nlg = AdvancedNLG()
        self.matcher = SemanticMatcher()
        self.conversation_memory = []
    
    def generate_response(self, question, content):
        """Generate natural, intelligent response"""
        if not content or len(content.strip()) < 20:
            return "I don't have enough content to answer that question. Could you upload or enter some text first?"
        
        # Find relevant content
        relevant_sentences, analysis = self.matcher.find_relevant_content(question, content)
        
        if not relevant_sentences or relevant_sentences[0]['score'] < 1:
            return self.generate_no_match_response(question, analysis)
        
        # Generate response based on question type
        if analysis['type'] == 'definition':
            response = self.generate_definition_response(relevant_sentences, analysis)
        elif analysis['type'] == 'explanation':
            response = self.generate_explanation_response(relevant_sentences, analysis)
        elif analysis['type'] == 'procedure':
            response = self.generate_procedure_response(relevant_sentences, analysis)
        elif analysis['type'] == 'comparison':
            response = self.generate_comparison_response(relevant_sentences, analysis)
        elif analysis['type'] == 'listing':
            response = self.generate_list_response(relevant_sentences, analysis)
        elif analysis['type'] == 'yes_no':
            response = self.generate_yes_no_response(relevant_sentences, analysis)
        else:
            response = self.generate_general_response(relevant_sentences, analysis)
        
        # Store in conversation memory
        self.conversation_memory.append({
            'question': question,
            'response': response,
            'timestamp': datetime.now()
        })
        
        return response
    
    def generate_definition_response(self, sentences, analysis):
        """Generate definition-style response"""
        intros = [
            "Let me explain what I found:",
            "Based on the document, here's the definition:",
            "From what I can see,",
            "According to the content,",
        ]
        
        main_sentence = self.nlg.paraphrase_intelligently(sentences[0]['text'])
        
        response = f"{random.choice(intros)} {main_sentence}"
        
        # Add supporting detail
        if len(sentences) > 1 and sentences[1]['score'] > 2:
            support = self.nlg.paraphrase_intelligently(sentences[1]['text'])
            response += f" {random.choice(self.nlg.transitions['elaboration'])} {support}"
        
        # Add engagement
        if random.random() > 0.5:
            response += "\n\n" + random.choice(self.nlg.engagers)
        
        return response
    
    def generate_explanation_response(self, sentences, analysis):
        """Generate explanatory response"""
        intros = [
            "Here's how this works:",
            "Let me break this down for you:",
            "The reason behind this is interesting:",
            "From what the document explains,",
        ]
        
        # Use top 3 sentences
        top_sentences = [self.nlg.paraphrase_intelligently(s['text']) for s in sentences[:3]]
        
        response = f"{random.choice(intros)} {top_sentences[0]}"
        
        if len(top_sentences) > 1:
            response += f" {random.choice(self.nlg.transitions['addition'])} {top_sentences[1]}"
        
        if len(top_sentences) > 2 and sentences[2]['score'] > 2:
            response += f" {random.choice(self.nlg.transitions['result'])} {top_sentences[2]}"
        
        response += "\n\nDoes this explanation make sense?"
        
        return response
    
    def generate_procedure_response(self, sentences, analysis):
        """Generate step-by-step response"""
        response = "Here's the process I found in the document:\n\n"
        
        # Look for numbered steps or sequential info
        steps = []
        for s in sentences[:5]:
            text = s['text']
            # Check if it's a step
            if re.search(r'\bfirst\b|\bsecond\b|\bstep\b|\bthen\b|\bnext\b', text.lower()):
                steps.append(self.nlg.paraphrase_intelligently(text))
        
        if steps:
            for i, step in enumerate(steps, 1):
                response += f"{i}. {step}\n\n"
        else:
            # Create a flowing explanation instead
            response = "Let me walk you through this:\n\n"
            for s in sentences[:3]:
                response += f"• {self.nlg.paraphrase_intelligently(s['text'])}\n\n"
        
        response += "Would you like me to clarify any of these points?"
        
        return response
    
    def generate_list_response(self, sentences, analysis):
        """Generate list-style response"""
        response = "Here's what I found:\n\n"
        
        items = []
        for s in sentences[:5]:
            text = s['text']
            # Look for list indicators
            if ',' in text or 'include' in text.lower():
                items.append(self.nlg.paraphrase_intelligently(text))
        
        if items:
            for item in items:
                response += f"• {item}\n\n"
        else:
            response = "Based on the content:\n\n"
            for s in sentences[:4]:
                response += f"• {self.nlg.paraphrase_intelligently(s['text'])}\n\n"
        
        return response.strip()
    
    def generate_yes_no_response(self, sentences, analysis):
        """Generate yes/no response with explanation"""
        # Determine yes or no based on content
        top_sentence = sentences[0]['text'].lower()
        
        # Look for affirmative or negative indicators
        positive_words = ['yes', 'correct', 'true', 'indeed', 'certainly', 'does', 'is', 'can', 'will']
        negative_words = ['no', 'not', 'never', 'cannot', 'won\'t', 'doesn\'t', 'isn\'t']
        
        pos_count = sum(1 for word in positive_words if word in top_sentence)
        neg_count = sum(1 for word in negative_words if word in top_sentence)
        
        if pos_count > neg_count:
            answer = "Yes"
        elif neg_count > pos_count:
            answer = "No"
        else:
            answer = "Based on what I found"
        
        main_text = self.nlg.paraphrase_intelligently(sentences[0]['text'])
        
        response = f"{answer}, {main_text}"
        
        # Add supporting evidence
        if len(sentences) > 1:
            support = self.nlg.paraphrase_intelligently(sentences[1]['text'])
            response += f" {random.choice(self.nlg.transitions['elaboration'])} {support}"
        
        return response
    
    def generate_general_response(self, sentences, analysis):
        """Generate general informative response"""
        # Use the NLG to create a natural paragraph
        sentence_texts = [s['text'] for s in sentences[:3]]
        paraphrased = [self.nlg.paraphrase_intelligently(s) for s in sentence_texts]
        
        response = self.nlg.create_human_paragraph(paraphrased, style='helpful')
        
        return response
    
    def generate_comparison_response(self, sentences, analysis):
        """Generate comparison response"""
        response = "Let me compare these for you:\n\n"
        
        # Try to find contrasting information
        for s in sentences[:4]:
            text = s['text']
            if 'differ' in text.lower() or 'while' in text.lower() or 'whereas' in text.lower():
                response += f"{self.nlg.paraphrase_intelligently(text)}\n\n"
        
        if response == "Let me compare these for you:\n\n":
            # No explicit comparison found, provide available info
            response = "Here's what I found about this:\n\n"
            for s in sentences[:3]:
                response += f"• {self.nlg.paraphrase_intelligently(s['text'])}\n\n"
        
        return response.strip()
    
    def generate_no_match_response(self, question, analysis):
        """Generate helpful response when no match found"""
        responses = [
            f"I couldn't find specific information about {', '.join(analysis['concepts']['keywords'][:2])} in the uploaded content. Could you rephrase your question or ask about something else?",
            f"Hmm, I don't see details about {', '.join(analysis['concepts']['keywords'][:2])} in the document. What else would you like to know?",
            f"I searched for information related to your question, but couldn't find a clear answer in the content. Try asking about a different topic!",
            "That's a great question, but I don't have enough relevant information in the current document to answer it properly. Could you ask something else?",
        ]
        
        return random.choice(responses)

class PPTContentGenerator:
    """Generates structured PPT content from document"""
    
    def __init__(self):
        pass
    
    def extract_key_points(self, text, num_points=5):
        """Extract key points from text"""
        sentences = [s.strip() for s in re.split(r'[.!?]+', text) if len(s.strip()) > 20]
        
        if not sentences:
            return []
        
        scored = []
        for i, sent in enumerate(sentences[:20]):
            score = len(sent.split()) * (1 - i * 0.05)
            scored.append({'text': sent, 'score': score})
        
        scored.sort(key=lambda x: x['score'], reverse=True)
        return [s['text'] for s in scored[:num_points]]
    
    def extract_topics(self, text):
        """Extract main topics from text"""
        paragraphs = [p.strip() for p in text.split('\n\n') if len(p.strip()) > 50]
        
        if not paragraphs:
            sentences = [s.strip() for s in re.split(r'[.!?]+', text) if len(s.strip()) > 30]
            paragraphs = []
            for i in range(0, len(sentences), 3):
                paragraphs.append(' '.join(sentences[i:i+3]))
        
        topics = []
        for para in paragraphs[:6]:
            first_sent = re.split(r'[.!?]+', para)[0].strip()
            words = first_sent.split()[:8]
            topic_title = ' '.join(words)
            
            points = self.extract_key_points(para, num_points=3)
            
            topics.append({
                'title': topic_title,
                'content': para[:300],
                'points': points[:3]
            })
        
        return topics
    
    def create_title_from_content(self, text):
        """Create a title from the content"""
        first_sentences = re.split(r'[.!?]+', text)[:3]
        common_words = {}
        
        for sent in first_sentences:
            words = re.findall(r'\b[A-Z][a-z]+\b', sent)
            for word in words:
                common_words[word] = common_words.get(word, 0) + 1
        
        if common_words:
            main_word = max(common_words, key=common_words.get)
            return f"{main_word} Overview"
        
        return "Content Overview"
    
    def generate_slide_content(self, content_text):
        """Generate 7-8 slides of content"""
        if not content_text or len(content_text) < 100:
            return None
        
        slides = []
        
        # Slide 1: Title Slide
        title = self.create_title_from_content(content_text)
        slides.append({
            'type': 'title',
            'title': title,
            'subtitle': f'Generated from document content\n{datetime.now().strftime("%B %d, %Y")}'
        })
        
        # Slide 2: Introduction/Overview
        intro_sentences = self.extract_key_points(content_text[:1000], num_points=4)
        slides.append({
            'type': 'content',
            'title': 'Introduction',
            'bullets': intro_sentences
        })
        
        # Slide 3-7: Topic slides
        topics = self.extract_topics(content_text)
        
        for topic in topics[:5]:
            slide_title = topic['title']
            if len(slide_title) > 60:
                slide_title = slide_title[:60] + "..."
            
            bullets = topic['points'] if topic['points'] else [topic['content'][:200]]
            
            slides.append({
                'type': 'content',
                'title': slide_title,
                'bullets': bullets
            })
        
        # Slide: Key Takeaways
        all_sentences = [s.strip() for s in re.split(r'[.!?]+', content_text) if len(s.strip()) > 30]
        key_takeaways = []
        
        for sent in all_sentences:
            if any(word in sent.lower() for word in ['important', 'key', 'essential', 'critical', 'significant', 'main', 'primary']):
                key_takeaways.append(sent)
                if len(key_takeaways) >= 4:
                    break
        
        if not key_takeaways:
            key_takeaways = self.extract_key_points(content_text, num_points=4)
        
        slides.append({
            'type': 'content',
            'title': 'Key Takeaways',
            'bullets': key_takeaways[:4]
        })
        
        # Final Slide: Conclusion
        slides.append({
            'type': 'conclusion',
            'title': 'Thank You',
            'subtitle': 'Questions?'
        })
        
        return slides

class ChatbotApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Advanced AI Chatbot (No ML Models)")
        self.root.geometry("1200x800")
        self.root.configure(bg='#f8f9fa')
        
        self.response_engine = AdvancedResponseEngine()
        self.ppt_generator = PPTContentGenerator()
        
        self.uploaded_content = ""
        self.current_summary = ""
        self.summary_visible = False
        self.message_history = []
        self.chat_sessions = []
        self.current_session_index = 0
        self.manual_input_visible = False
        self.session_counter = 1
        
        self.create_new_session()
        self.create_widgets()
        
        self.add_bot_message("Hey there! 👋 I'm your AI assistant. I can understand your questions and give natural, helpful responses based on any document you upload. I use advanced algorithms (not machine learning!) to truly understand what you're asking. Upload a file or enter some text, and let's chat!")
    
    def create_widgets(self):
        main_frame = tk.Frame(self.root, bg='#f8f9fa')
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        self.create_sidebar(main_frame)
        self.create_main_content(main_frame)
    
    def create_sidebar(self, parent):
        sidebar_frame = tk.Frame(parent, bg='white', width=280, relief=tk.RAISED, bd=1)
        sidebar_frame.pack(side=tk.LEFT, fill=tk.Y, padx=(0, 1))
        sidebar_frame.pack_propagate(False)
        
        header_frame = tk.Frame(sidebar_frame, bg='#f8f9fa', height=80)
        header_frame.pack(fill=tk.X, padx=10, pady=10)
        header_frame.pack_propagate(False)
        
        new_chat_btn = tk.Button(
            header_frame, 
            text="+ New Chat",
            command=self.start_new_chat,
            bg='#007bff',
            fg='white',
            font=('Arial', 10, 'bold'),
            relief=tk.FLAT,
            cursor='hand2',
            pady=8
        )
        new_chat_btn.pack(fill=tk.X, pady=5)
        
        history_label = tk.Label(sidebar_frame, text="Chat History", font=('Arial', 9, 'bold'), bg='white', fg='#666')
        history_label.pack(anchor=tk.W, padx=10, pady=(10, 5))
        
        self.history_listbox = tk.Listbox(
            sidebar_frame,
            font=('Arial', 9),
            bg='#f8f9fa',
            selectbackground='#007bff',
            selectforeground='white',
            relief=tk.FLAT,
            bd=0,
            highlightthickness=0
        )
        self.history_listbox.pack(fill=tk.BOTH, expand=True, padx=10, pady=(0, 10))
        
        self.update_history_display()
        self.history_listbox.bind('<<ListboxSelect>>', self.on_history_select)
    
    def create_main_content(self, parent):
        main_content = tk.Frame(parent, bg='#f8f9fa')
        main_content.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True)
        
        self.create_upload_section(main_content)
        self.create_summary_section(main_content)
        self.create_chat_area(main_content)
        self.create_input_section(main_content)
    
    def create_upload_section(self, parent):
        upload_main_frame = tk.Frame(parent, bg='white', relief=tk.RAISED, bd=1)
        upload_main_frame.pack(fill=tk.X, padx=20, pady=(20, 10))
        
        upload_frame = tk.Frame(upload_main_frame, bg='white')
        upload_frame.pack(fill=tk.X, padx=15, pady=15)
        
        label = tk.Label(upload_frame, text="📄 Upload Document or Enter Text", 
                        font=('Arial', 11, 'bold'), bg='white', fg='#333')
        label.pack(anchor=tk.W, pady=(0, 10))
        
        button_frame = tk.Frame(upload_frame, bg='white')
        button_frame.pack(fill=tk.X)
        
        upload_btn = tk.Button(
            button_frame,
            text="📁 Upload File",
            command=self.upload_file,
            bg='#28a745',
            fg='white',
            font=('Arial', 9, 'bold'),
            relief=tk.FLAT,
            cursor='hand2',
            padx=15,
            pady=8
        )
        upload_btn.pack(side=tk.LEFT, padx=(0, 10))
        
        manual_input_btn = tk.Button(
            button_frame,
            text="✏️ Enter Text Manually",
            command=self.toggle_manual_input,
            bg='#6c757d',
            fg='white',
            font=('Arial', 9, 'bold'),
            relief=tk.FLAT,
            cursor='hand2',
            padx=15,
            pady=8
        )
        manual_input_btn.pack(side=tk.LEFT, padx=(0, 10))
        
        if PPTX_WRITE_AVAILABLE:
            ppt_btn = tk.Button(
                button_frame,
                text="📊 Generate PPT",
                command=self.generate_ppt,
                bg='#fd7e14',
                fg='white',
                font=('Arial', 9, 'bold'),
                relief=tk.FLAT,
                cursor='hand2',
                padx=15,
                pady=8
            )
            ppt_btn.pack(side=tk.LEFT)
        
        self.manual_input_frame = tk.Frame(upload_main_frame, bg='white')
        
        manual_label = tk.Label(self.manual_input_frame, text="Enter your text:", 
                               font=('Arial', 9), bg='white', fg='#666')
        manual_label.pack(anchor=tk.W, padx=15, pady=(10, 5))
        
        self.manual_text = scrolledtext.ScrolledText(
            self.manual_input_frame,
            height=6,
            font=('Arial', 9),
            wrap=tk.WORD,
            relief=tk.FLAT,
            bg='#f8f9fa',
            fg='#333'
        )
        self.manual_text.pack(fill=tk.BOTH, padx=15, pady=(0, 10))
        
        manual_btn_frame = tk.Frame(self.manual_input_frame, bg='white')
        manual_btn_frame.pack(fill=tk.X, padx=15, pady=(0, 15))
        
        save_manual_btn = tk.Button(
            manual_btn_frame,
            text="Save Text",
            command=self.save_manual_input,
            bg='#007bff',
            fg='white',
            font=('Arial', 9, 'bold'),
            relief=tk.FLAT,
            cursor='hand2',
            padx=20,
            pady=6
        )
        save_manual_btn.pack(side=tk.LEFT, padx=(0, 10))
        
        cancel_manual_btn = tk.Button(
            manual_btn_frame,
            text="Cancel",
            command=self.toggle_manual_input,
            bg='#dc3545',
            fg='white',
            font=('Arial', 9, 'bold'),
            relief=tk.FLAT,
            cursor='hand2',
            padx=20,
            pady=6
        )
        cancel_manual_btn.pack(side=tk.LEFT)
        
        self.file_info_label = tk.Label(
            upload_main_frame,
            text="",
            font=('Arial', 8),
            bg='white',
            fg='#666',
            anchor=tk.W
        )
        self.file_info_label.pack(fill=tk.X, padx=15, pady=(0, 10))
    
    def create_summary_section(self, parent):
        self.summary_frame = tk.Frame(parent, bg='white', relief=tk.RAISED, bd=1)
        
        summary_header = tk.Frame(self.summary_frame, bg='#e9ecef')
        summary_header.pack(fill=tk.X)
        
        summary_title = tk.Label(
            summary_header,
            text="📋 Document Summary",
            font=('Arial', 10, 'bold'),
            bg='#e9ecef',
            fg='#333'
        )
        summary_title.pack(side=tk.LEFT, padx=15, pady=10)
        
        self.summary_text = scrolledtext.ScrolledText(
            self.summary_frame,
            height=5,
            font=('Arial', 9),
            wrap=tk.WORD,
            relief=tk.FLAT,
            bg='#f8f9fa',
            fg='#333',
            state=tk.DISABLED
        )
        self.summary_text.pack(fill=tk.BOTH, padx=15, pady=(0, 15))
    
    def create_chat_area(self, parent):
        chat_frame = tk.Frame(parent, bg='white', relief=tk.RAISED, bd=1)
        chat_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        self.chat_display = scrolledtext.ScrolledText(
            chat_frame,
            font=('Arial', 10),
            wrap=tk.WORD,
            relief=tk.FLAT,
            bg='#ffffff',
            fg='#333',
            state=tk.DISABLED,
            padx=15,
            pady=15
        )
        self.chat_display.pack(fill=tk.BOTH, expand=True)
        
        self.chat_display.tag_config('user', foreground='#007bff', font=('Arial', 10, 'bold'))
        self.chat_display.tag_config('bot', foreground='#28a745', font=('Arial', 10, 'bold'))
        self.chat_display.tag_config('timestamp', foreground='#999', font=('Arial', 8))
        self.chat_display.tag_config('message', foreground='#333', font=('Arial', 10))
    
    def create_input_section(self, parent):
        input_frame = tk.Frame(parent, bg='white', relief=tk.RAISED, bd=1)
        input_frame.pack(fill=tk.X, padx=20, pady=(0, 20))
        
        input_container = tk.Frame(input_frame, bg='white')
        input_container.pack(fill=tk.X, padx=15, pady=15)
        
        self.user_input = tk.Text(
            input_container,
            height=3,
            font=('Arial', 10),
            wrap=tk.WORD,
            relief=tk.FLAT,
            bg='#f8f9fa',
            fg='#333',
            padx=10,
            pady=10
        )
        self.user_input.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(0, 10))
        self.user_input.bind('<Return>', self.handle_enter)
        
        send_btn = tk.Button(
            input_container,
            text="Send ➤",
            command=self.send_message,
            bg='#007bff',
            fg='white',
            font=('Arial', 10, 'bold'),
            relief=tk.FLAT,
            cursor='hand2',
            padx=20,
            pady=10
        )
        send_btn.pack(side=tk.RIGHT)
    
    def toggle_manual_input(self):
        if self.manual_input_visible:
            self.manual_input_frame.pack_forget()
            self.manual_input_visible = False
        else:
            self.manual_input_frame.pack(fill=tk.BOTH, padx=0, pady=(10, 0))
            self.manual_input_visible = True
    
    def save_manual_input(self):
        text = self.manual_text.get('1.0', tk.END).strip()
        if text:
            self.uploaded_content = text
            self.chat_sessions[self.current_session_index]['content'] = text
            self.file_info_label.config(text=f"✓ Manual text loaded ({len(text)} characters)")
            self.generate_summary(text)
            self.toggle_manual_input()
            self.add_bot_message("Got it! I've processed your text. Ask me anything about it!")
        else:
            messagebox.showwarning("Empty Text", "Please enter some text before saving.")
    
    def upload_file(self):
        file_path = filedialog.askopenfilename(
            title="Select a file",
            filetypes=[
                ("Text files", "*.txt"),
                ("PDF files", "*.pdf"),
                ("PowerPoint", "*.pptx"),
                ("All files", "*.*")
            ]
        )
        
        if file_path:
            self.process_file(file_path)
    
    def process_file(self, file_path):
        filename = os.path.basename(file_path)
        extension = os.path.splitext(filename)[1].lower()
        
        try:
            if extension == '.txt':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            
            elif extension == '.pdf' and PDF_AVAILABLE:
                content = self.extract_pdf_text(file_path)
            
            elif extension == '.pptx' and PPTX_AVAILABLE:
                content = self.extract_pptx_text(file_path)
            
            else:
                messagebox.showerror("Unsupported Format", 
                                   f"Cannot process {extension} files. Please upload .txt, .pdf, or .pptx files.")
                return
            
            if content and len(content.strip()) > 0:
                self.uploaded_content = content
                self.chat_sessions[self.current_session_index]['content'] = content
                self.file_info_label.config(text=f"✓ {filename} loaded successfully ({len(content)} characters)")
                self.generate_summary(content)
                self.add_bot_message(f"Great! I've analyzed '{filename}'. What would you like to know?")
            else:
                messagebox.showwarning("Empty File", "The file appears to be empty or unreadable.")
        
        except Exception as e:
            messagebox.showerror("Error", f"Error processing file: {str(e)}")
    
    def extract_pdf_text(self, file_path):
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            raise Exception(f"PDF extraction error: {str(e)}")
        return text
    
    def extract_pptx_text(self, file_path):
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            raise Exception(f"PPTX extraction error: {str(e)}")
        return text
    
    def generate_summary(self, content):
        sentences = [s.strip() for s in re.split(r'[.!?]+', content) if len(s.strip()) > 20]
        
        if len(sentences) < 3:
            summary = content[:500]
        else:
            scored = []
            for i, sent in enumerate(sentences[:15]):
                score = len(sent.split()) * (1 - i * 0.05)
                scored.append({'text': sent, 'score': score})
            
            scored.sort(key=lambda x: x['score'], reverse=True)
            top_sentences = [s['text'] for s in scored[:3]]
            summary = '. '.join(top_sentences) + '.'
        
        self.current_summary = summary
        
        if not self.summary_visible:
            self.summary_frame.pack(fill=tk.X, padx=20, pady=(0, 10))
            self.summary_visible = True
        
        self.summary_text.config(state=tk.NORMAL)
        self.summary_text.delete('1.0', tk.END)
        self.summary_text.insert('1.0', summary)
        self.summary_text.config(state=tk.DISABLED)
    
    def generate_ppt(self):
        if not self.uploaded_content:
            messagebox.showwarning("No Content", "Please upload or enter some content first.")
            return
        
        if not PPTX_WRITE_AVAILABLE:
            messagebox.showerror("Not Available", "PowerPoint generation requires python-pptx library.")
            return
        
        self.add_bot_message("Generating PowerPoint presentation... ⏳")
        
        def generate_thread():
            try:
                slides_content = self.ppt_generator.generate_slide_content(self.uploaded_content)
                
                if not slides_content:
                    self.root.after(0, lambda: messagebox.showwarning("Insufficient Content", 
                                    "Need more content to generate a meaningful presentation."))
                    return
                
                prs = Presentation()
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
                
                for slide_data in slides_content:
                    if slide_data['type'] == 'title':
                        slide = prs.slides.add_slide(prs.slide_layouts[0])
                        slide.shapes.title.text = slide_data['title']
                        if len(slide.shapes) > 1:
                            slide.placeholders[1].text = slide_data['subtitle']
                    
                    elif slide_data['type'] == 'content':
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        slide.shapes.title.text = slide_data['title']
                        
                        body_shape = slide.shapes.placeholders[1]
                        tf = body_shape.text_frame
                        
                        for bullet in slide_data['bullets']:
                            p = tf.add_paragraph()
                            p.text = bullet[:150]
                            p.level = 0
                    
                    elif slide_data['type'] == 'conclusion':
                        slide = prs.slides.add_slide(prs.slide_layouts[0])
                        slide.shapes.title.text = slide_data['title']
                        if len(slide.shapes) > 1:
                            slide.placeholders[1].text = slide_data['subtitle']
                
                save_path = filedialog.asksaveasfilename(
                    defaultextension=".pptx",
                    filetypes=[("PowerPoint", "*.pptx")],
                    initialfile="generated_presentation.pptx"
                )
                
                if save_path:
                    prs.save(save_path)
                    self.root.after(0, lambda: self.add_bot_message(
                        f"✅ PowerPoint generated successfully! Saved to: {os.path.basename(save_path)}"))
                    self.root.after(0, lambda: messagebox.showinfo("Success", 
                                    f"Presentation created with {len(slides_content)} slides!"))
            
            except Exception as e:
                self.root.after(0, lambda: messagebox.showerror("Error", f"Error generating PPT: {str(e)}"))
                self.root.after(0, lambda: self.add_bot_message("Sorry, there was an error generating the PPT."))
        
        thread = threading.Thread(target=generate_thread)
        thread.daemon = True
        thread.start()
    
    def handle_enter(self, event):
        if event.state & 0x1:  # Shift+Enter
            return
        else:
            self.send_message()
            return 'break'
    
    def send_message(self):
        user_message = self.user_input.get('1.0', tk.END).strip()
        
        if not user_message:
            return
        
        self.user_input.delete('1.0', tk.END)
        
        self.add_user_message(user_message)
        
        self.root.after(100, lambda: self.generate_bot_response(user_message))
    
    def add_user_message(self, message):
        timestamp = datetime.now().strftime("%H:%M")
        
        self.chat_display.config(state=tk.NORMAL)
        
        self.chat_display.insert(tk.END, "You", 'user')
        self.chat_display.insert(tk.END, f" • {timestamp}\n", 'timestamp')
        self.chat_display.insert(tk.END, f"{message}\n\n", 'message')
        
        self.chat_display.config(state=tk.DISABLED)
        self.chat_display.see(tk.END)
        
        self.message_history.append({
            'role': 'user',
            'message': message,
            'timestamp': timestamp
        })
        
        self.chat_sessions[self.current_session_index]['messages'].append({
            'role': 'user',
            'message': message,
            'timestamp': timestamp
        })
    
    def add_bot_message(self, message):
        timestamp = datetime.now().strftime("%H:%M")
        
        self.chat_display.config(state=tk.NORMAL)
        
        self.chat_display.insert(tk.END, "Assistant", 'bot')
        self.chat_display.insert(tk.END, f" • {timestamp}\n", 'timestamp')
        self.chat_display.insert(tk.END, f"{message}\n\n", 'message')
        
        self.chat_display.config(state=tk.DISABLED)
        self.chat_display.see(tk.END)
        
        self.message_history.append({
            'role': 'bot',
            'message': message,
            'timestamp': timestamp
        })
        
        self.chat_sessions[self.current_session_index]['messages'].append({
            'role': 'bot',
            'message': message,
            'timestamp': timestamp
        })
    
    def generate_bot_response(self, user_message):
        response = self.response_engine.generate_response(user_message, self.uploaded_content)
        self.add_bot_message(response)
    
    def create_new_session(self):
        session = {
            'id': self.session_counter,
            'name': f"Chat {self.session_counter}",
            'messages': [],
            'content': "",
            'created': datetime.now()
        }
        self.chat_sessions.append(session)
        self.session_counter += 1
        return len(self.chat_sessions) - 1
    
    def start_new_chat(self):
        new_index = self.create_new_session()
        self.current_session_index = new_index
        
        self.uploaded_content = ""
        self.current_summary = ""
        self.message_history = []
        
        self.chat_display.config(state=tk.NORMAL)
        self.chat_display.delete('1.0', tk.END)
        self.chat_display.config(state=tk.DISABLED)
        
        if self.summary_visible:
            self.summary_frame.pack_forget()
            self.summary_visible = False
        
        self.file_info_label.config(text="")
        
        self.update_history_display()
        
        self.add_bot_message("New chat started! Upload a document or enter text to begin.")
    
    def update_history_display(self):
        self.history_listbox.delete(0, tk.END)
        
        for i, session in enumerate(self.chat_sessions):
            display_text = f"  {session['name']}"
            if i == self.current_session_index:
                display_text = f"▶ {session['name']}"
            self.history_listbox.insert(tk.END, display_text)
    
    def on_history_select(self, event):
        selection = self.history_listbox.curselection()
        if selection:
            index = selection[0]
            if index != self.current_session_index:
                self.load_session(index)
    
    def load_session(self, index):
        self.current_session_index = index
        session = self.chat_sessions[index]
        
        self.uploaded_content = session['content']
        self.message_history = session['messages'].copy()
        
        self.chat_display.config(state=tk.NORMAL)
        self.chat_display.delete('1.0', tk.END)
        self.chat_display.config(state=tk.DISABLED)
        
        for msg in session['messages']:
            if msg['role'] == 'user':
                self.chat_display.config(state=tk.NORMAL)
                self.chat_display.insert(tk.END, "You", 'user')
                self.chat_display.insert(tk.END, f" • {msg['timestamp']}\n", 'timestamp')
                self.chat_display.insert(tk.END, f"{msg['message']}\n\n", 'message')
                self.chat_display.config(state=tk.DISABLED)
            else:
                self.chat_display.config(state=tk.NORMAL)
                self.chat_display.insert(tk.END, "Assistant", 'bot')
                self.chat_display.insert(tk.END, f" • {msg['timestamp']}\n", 'timestamp')
                self.chat_display.insert(tk.END, f"{msg['message']}\n\n", 'message')
                self.chat_display.config(state=tk.DISABLED)
        
        if self.uploaded_content:
            self.generate_summary(self.uploaded_content)
            self.file_info_label.config(text=f"✓ Content loaded ({len(self.uploaded_content)} characters)")
        else:
            if self.summary_visible:
                self.summary_frame.pack_forget()
                self.summary_visible = False
            self.file_info_label.config(text="")
        
        self.update_history_display()
        self.chat_display.see(tk.END)

def main():
    root = tk.Tk()
    app = ChatbotApp(root)
    root.mainloop()

if __name__ == "__main__":
    main()